"""
AgriSmart PPT Generator
Creates a professional PowerPoint presentation for the AgriSmart AI-Powered Vegetable Marketplace.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── Color Palette ────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x16, 0xa3, 0x4a)   # #16a34a – brand green
GREEN_LIGHT = RGBColor(0x4a, 0xde, 0x80)   # #4ade80
BLUE        = RGBColor(0x25, 0x63, 0xeb)   # #2563eb
BLUE_LIGHT  = RGBColor(0x93, 0xc5, 0xfd)   # #93c5fd
AMBER       = RGBColor(0xf5, 0x9e, 0x0b)   # #f59e0b
PURPLE      = RGBColor(0xa8, 0x55, 0xf7)   # #a855f7
RED         = RGBColor(0xef, 0x44, 0x44)   # #ef4444
DARK_BG     = RGBColor(0x0f, 0x17, 0x2a)   # #0f172a – slide background
CARD_BG     = RGBColor(0x1e, 0x29, 0x3b)   # #1e293b
TEXT_MAIN   = RGBColor(0xe2, 0xe8, 0xf0)   # #e2e8f0
TEXT_SUB    = RGBColor(0x94, 0xa3, 0xb8)   # #94a3b8
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, alpha=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(14), bold=False, color=TEXT_MAIN,
             align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_size=Pt(13), bold_first=False, color=TEXT_SUB, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = (i == 0 and bold_first)
        run.font.color.rgb = color


def add_accent_line(slide, x, y, w, color=GREEN, height=Pt(3)):
    shape = slide.shapes.add_shape(1, x, y, w, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_header(slide, title, subtitle=None, accent_color=GREEN):
    set_bg(slide, DARK_BG)
    # Top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), accent_color)
    # Title
    add_text(slide, title,
             Inches(0.5), Inches(0.18),
             Inches(12), Inches(0.75),
             font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.88),
                 Inches(11), Inches(0.42),
                 font_size=Pt(14), color=TEXT_SUB)
    # Divider
    add_accent_line(slide, Inches(0.5), Inches(1.28), Inches(12.3), accent_color)


def add_card(slide, x, y, w, h, border_color=GREEN):
    add_rect(slide, x, y, w, h, CARD_BG, line_color=border_color, line_width=Pt(1.5))


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — Title / Cover
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, DARK_BG)

# Full-width green gradient bar at top
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), GREEN)

# Large centered card
add_rect(slide,
         Inches(1.5), Inches(0.8),
         Inches(10.3), Inches(5.9),
         CARD_BG,
         line_color=GREEN, line_width=Pt(1.5))

# Accent stripe inside card
add_rect(slide, Inches(1.5), Inches(0.8), Inches(0.18), Inches(5.9), GREEN)

# Big title
add_text(slide, "AgriSmart",
         Inches(2.0), Inches(1.1),
         Inches(9.5), Inches(1.2),
         font_size=Pt(64), bold=True, color=GREEN_LIGHT,
         align=PP_ALIGN.CENTER)

add_text(slide, "AI-Powered Vegetable Marketplace",
         Inches(2.0), Inches(2.2),
         Inches(9.5), Inches(0.7),
         font_size=Pt(24), bold=False, color=WHITE,
         align=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(3.5), Inches(2.95), Inches(6.3), GREEN)

add_text(slide, "Connecting Farmers & Consumers across Andhra Pradesh\nwith Real-time AI Price Predictions & Live Location Mapping",
         Inches(2.0), Inches(3.1),
         Inches(9.5), Inches(1.0),
         font_size=Pt(15), color=TEXT_SUB,
         align=PP_ALIGN.CENTER)

# Stat pills
stats = [
    ("200+ Farmers", GREEN_LIGHT),
    ("13 Districts", BLUE_LIGHT),
    ("AI Price ML", AMBER),
    ("Live Map", PURPLE),
]
pill_x = Inches(1.8)
for label, col in stats:
    add_rect(slide, pill_x, Inches(4.4), Inches(2.1), Inches(0.52), DARK_BG,
             line_color=col, line_width=Pt(1))
    add_text(slide, label, pill_x + Inches(0.1), Inches(4.43), Inches(1.9), Inches(0.48),
             font_size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    pill_x += Inches(2.25)

add_text(slide, "Digital Transparency Initiative · Andhra Pradesh Agricultural Tech",
         Inches(1.5), Inches(5.95),
         Inches(10.3), Inches(0.4),
         font_size=Pt(10), color=TEXT_SUB,
         align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem Statement
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide,
             "The Problem",
             "Challenges Facing Farmers & Consumers in Andhra Pradesh",
             RED)

problems = [
    ("Price Opacity",
     "Farmers receive low farm-gate prices set by middlemen while consumers pay high market prices. There is no transparent price discovery mechanism."),
    ("Middlemen Exploitation",
     "Agricultural supply chains involve 3–5 intermediaries, each adding margin. Farmers may earn only 20–30% of the final retail price."),
    ("No Direct Channel",
     "Farmers and consumers have no direct platform to connect. Produce often spoils before it reaches the right buyer at the right location."),
    ("Unpredictable Price Volatility",
     "Seasonal fluctuations in vegetable prices are not communicated in advance, making it hard for farmers to plan harvests and for consumers to budget."),
    ("Geographic Information Gap",
     "Consumers don't know which farmer in their district has fresh produce available today, leading to food waste and supply-demand mismatches."),
]

colors = [RED, AMBER, PURPLE, BLUE, GREEN]
for i, (title, body) in enumerate(problems):
    row = i // 3
    col = i % 3
    if i < 3:
        cx = Inches(0.35 + col * 4.32)
        cy = Inches(1.55)
        cw = Inches(4.1)
        ch = Inches(2.35)
    else:
        cx = Inches(1.75 + (i - 3) * 4.32)
        cy = Inches(4.05)
        cw = Inches(4.1)
        ch = Inches(2.45)

    add_card(slide, cx, cy, cw, ch, colors[i])
    add_rect(slide, cx, cy, cw, Inches(0.07), colors[i])
    add_text(slide, title,
             cx + Inches(0.15), cy + Inches(0.12),
             cw - Inches(0.2), Inches(0.4),
             font_size=Pt(13), bold=True, color=colors[i])
    add_text(slide, body,
             cx + Inches(0.15), cy + Inches(0.52),
             cw - Inches(0.25), ch - Inches(0.7),
             font_size=Pt(10.5), color=TEXT_SUB)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Solution Overview
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "Our Solution — AgriSmart Platform",
             "One Unified Platform. AI + Maps + Direct Farmer-Consumer Connection.", GREEN)

# Left column — key pillars
pillars = [
    (GREEN,  "🌾", "Farmer Portal",     "Login, list your vegetables with quantity & district, manage availability, get AI price guidance."),
    (BLUE,   "🛒", "Consumer Portal",   "Search by vegetable or district, view live prices, find nearest farmers on an interactive map."),
    (AMBER,  "🤖", "AI Price Engine",   "Random Forest ML model trained on 3-year historical data predicts real-time market prices."),
    (PURPLE, "🗺️", "Live Map Search",   "Leaflet-powered map displays farmer pin-drops with exact GPS coordinates per district."),
]

for i, (col, icon, title, body) in enumerate(pillars):
    row = i // 2
    c   = i %  2
    cx = Inches(0.4 + c * 6.35)
    cy = Inches(1.55 + row * 2.7)
    cw = Inches(6.1)
    ch = Inches(2.45)
    add_card(slide, cx, cy, cw, ch, col)
    add_rect(slide, cx, cy, Inches(0.14), ch, col)
    add_text(slide, f"{icon}  {title}",
             cx + Inches(0.25), cy + Inches(0.15),
             cw - Inches(0.3), Inches(0.45),
             font_size=Pt(14), bold=True, color=col)
    add_text(slide, body,
             cx + Inches(0.25), cy + Inches(0.6),
             cw - Inches(0.35), ch - Inches(0.8),
             font_size=Pt(11.5), color=TEXT_SUB)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Tech Stack
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "Technology Stack",
             "Modern, Scalable, Open-Source Architecture", BLUE)

layers = [
    ("FRONTEND",  BLUE,   [
        ("React 18 + TypeScript",   "Type-safe, component-based UI"),
        ("Vite",                    "Ultra-fast HMR dev server"),
        ("Tailwind CSS",            "Utility-first, responsive design"),
        ("Leaflet.js",              "Interactive map with geo-pins"),
        ("Lucide Icons",            "Crisp SVG icon library"),
    ]),
    ("BACKEND",   GREEN,  [
        ("Python 3.11 + FastAPI",   "Async REST API, auto-docs via Swagger"),
        ("Pandas + NumPy",          "Data wrangling & analytics"),
        ("Scikit-learn",            "Random Forest ML model"),
        ("Joblib",                  "Model serialization & fast loading"),
    ]),
    ("AUTH & DATA", AMBER, [
        ("JWT + bcrypt",            "Secure token auth with hashed passwords"),
        ("File-based JSON store",   "Lightweight users.json & crops.json"),
        ("CSV Datasets",            "3-year price history, weather, soil, production data"),
    ]),
]

col_x = [Inches(0.35), Inches(4.7), Inches(8.85)]
for i, (layer_title, col, items) in enumerate(layers):
    cx = col_x[i]
    cy = Inches(1.5)
    cw = Inches(4.15)
    ch = Inches(5.7)
    add_card(slide, cx, cy, cw, ch, col)
    add_rect(slide, cx, cy, cw, Inches(0.55), col)
    add_text(slide, layer_title,
             cx + Inches(0.1), cy + Inches(0.08),
             cw - Inches(0.15), Inches(0.42),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (tech, desc) in enumerate(items):
        iy = cy + Inches(0.7 + j * 0.95)
        add_text(slide, f"▸ {tech}",
                 cx + Inches(0.15), iy,
                 cw - Inches(0.25), Inches(0.35),
                 font_size=Pt(12), bold=True, color=col)
        add_text(slide, desc,
                 cx + Inches(0.2), iy + Inches(0.34),
                 cw - Inches(0.3), Inches(0.42),
                 font_size=Pt(10), color=TEXT_SUB)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — AI / ML Model
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "AI Price Prediction Engine",
             "Random Forest ML Model — Trained on 3-Year Andhra Pradesh Market Data", AMBER)

# Left: Model info
add_card(slide, Inches(0.35), Inches(1.55), Inches(6.0), Inches(5.7), AMBER)
add_rect(slide, Inches(0.35), Inches(1.55), Inches(6.0), Inches(0.55), AMBER)
add_text(slide, "HOW IT WORKS",
         Inches(0.45), Inches(1.6),
         Inches(5.8), Inches(0.45),
         font_size=Pt(12.5), bold=True, color=WHITE)

steps = [
    ("1. Data Ingestion",
     "Loads prices_3year_dataset.csv (~1.8 MB) with historical vegetable prices by district, month, and year."),
    ("2. Feature Engineering",
     "Features: vegetable, district, month, year, avg_yield (kg/ha), avg_area (ha). Target: avg_price (₹/kg)."),
    ("3. Model Training",
     "Scikit-learn RandomForestRegressor trained on all data. Model & encoders serialized with joblib."),
    ("4. Auto-Load on Startup",
     "FastAPI lifespan checks for saved model. If absent, trains automatically — zero manual intervention."),
    ("5. Real-time Inference",
     "GET /api/predict-price returns price prediction within milliseconds per query."),
]

for i, (title, body) in enumerate(steps):
    iy = Inches(2.2 + i * 0.98)
    add_text(slide, title,
             Inches(0.55), iy,
             Inches(5.7), Inches(0.35),
             font_size=Pt(12), bold=True, color=AMBER)
    add_text(slide, body,
             Inches(0.6), iy + Inches(0.34),
             Inches(5.55), Inches(0.5),
             font_size=Pt(10.5), color=TEXT_SUB)

# Right: Features used + datasets
add_card(slide, Inches(6.7), Inches(1.55), Inches(6.28), Inches(2.65), GREEN)
add_rect(slide, Inches(6.7), Inches(1.55), Inches(6.28), Inches(0.5), GREEN)
add_text(slide, "DATASETS USED",
         Inches(6.8), Inches(1.6),
         Inches(6.0), Inches(0.4),
         font_size=Pt(12.5), bold=True, color=WHITE)

datasets = [
    ("prices_3year_dataset.csv", "Historical market prices — 1.8 MB"),
    ("veg_production_dataset.csv", "Yield & area data per vegetable/district"),
    ("weather_dataset.csv", "Seasonal weather context"),
    ("soil_dataset.csv", "Soil quality indicators"),
    ("districts_dataset.csv", "Lat/Lon GPS coordinates for all 13 AP districts"),
]
for i, (name, desc) in enumerate(datasets):
    iy = Inches(2.1 + i * 0.44)
    add_text(slide, f"📁 {name}",
             Inches(6.85), iy,
             Inches(6.0), Inches(0.28),
             font_size=Pt(10.5), bold=True, color=GREEN_LIGHT)
    add_text(slide, f"       {desc}",
             Inches(6.85), iy + Inches(0.25),
             Inches(6.0), Inches(0.22),
             font_size=Pt(9.5), color=TEXT_SUB)

add_card(slide, Inches(6.7), Inches(4.4), Inches(6.28), Inches(2.8), BLUE)
add_rect(slide, Inches(6.7), Inches(4.4), Inches(6.28), Inches(0.5), BLUE)
add_text(slide, "KEY METRICS",
         Inches(6.8), Inches(4.45),
         Inches(6.0), Inches(0.4),
         font_size=Pt(12.5), bold=True, color=WHITE)

metrics = [
    ("Algorithm", "Random Forest Regressor (Scikit-learn)"),
    ("Target Variable", "avg_price (₹ per kg)"),
    ("Evaluation Metric", "Mean Absolute Error (MAE)"),
    ("Inference Endpoint", "GET /api/predict-price"),
    ("Recommendations Endpoint", "GET /api/recommendations?district=Guntur"),
    ("Voice Command NLP", "POST /api/voice-command (regex-based intent parser)"),
]
for i, (k, v) in enumerate(metrics):
    iy = Inches(5.0 + i * 0.35)
    add_text(slide, f"  {k}: ",
             Inches(6.85), iy,
             Inches(2.3), Inches(0.32),
             font_size=Pt(10), bold=True, color=BLUE_LIGHT)
    add_text(slide, v,
             Inches(9.05), iy,
             Inches(3.75), Inches(0.32),
             font_size=Pt(10), color=TEXT_SUB)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Key Features
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "Key Features",
             "Built for Real-World Agricultural Transparency", GREEN)

features = [
    (GREEN,  "🌾", "Farmer Listing Portal",
     "Farmers register, login (JWT), and list vegetables with name, quantity (kg), district, and location. Each listing is pinned on the live map."),
    (BLUE,   "🗺️", "Interactive Map Search",
     "Consumers search by vegetable or district. Leaflet.js renders markers for every farmer & market with click-to-call contact cards."),
    (AMBER,  "🤖", "AI Price Predictions",
     "Users get real-time price forecasts for any vegetable+district+month combination. Recommendations surface high-yield, low-competition crops."),
    (PURPLE, "🎤", "Voice Commands",
     "NLP-lite voice command parser handles intents: add crop, get price, navigate pages. Works without any external API."),
    (GREEN_LIGHT, "📊", "System & Model Settings",
     "Admin Settings page shows MAE, training timestamp, model metadata, and allows triggering re-training from the UI."),
    (BLUE_LIGHT,  "🔐", "Secure Authentication",
     "Separate login portals for Farmers and Consumers. bcrypt password hashing. JWT tokens for session management."),
]

for i, (col, icon, title, body) in enumerate(features):
    row = i // 3
    c   = i %  3
    cx = Inches(0.35 + c * 4.32)
    cy = Inches(1.55 + row * 2.7)
    cw = Inches(4.1)
    ch = Inches(2.5)
    add_card(slide, cx, cy, cw, ch, col)
    add_rect(slide, cx, cy, Inches(0.12), ch, col)
    add_text(slide, f"{icon} {title}",
             cx + Inches(0.23), cy + Inches(0.15),
             cw - Inches(0.3), Inches(0.45),
             font_size=Pt(13), bold=True, color=col)
    add_text(slide, body,
             cx + Inches(0.23), cy + Inches(0.62),
             cw - Inches(0.35), ch - Inches(0.8),
             font_size=Pt(10.5), color=TEXT_SUB)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — System Architecture
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "System Architecture",
             "Full-Stack Flow: React Frontend → FastAPI Backend → ML Engine", PURPLE)

# Architecture flow diagram (text-based)
boxes = [
    (Inches(0.4),  Inches(2.5),  Inches(2.4), Inches(3.5), BLUE,   "FRONTEND\n(React + Vite)\nPort 5173",
     ["• Landing Page", "• Farmer Login/Portal", "• Consumer Search", "• Map (Leaflet)", "• Settings"]),
    (Inches(3.2),  Inches(2.5),  Inches(2.4), Inches(3.5), GREEN,  "BACKEND\n(FastAPI)\nPort 8000",
     ["• /api/vegetables", "• /api/predict-price", "• /api/recommendations", "• /api/voice-command", "• /auth/*"]),
    (Inches(6.0),  Inches(2.5),  Inches(2.4), Inches(3.5), AMBER,  "ML ENGINE\n(Scikit-learn)",
     ["• Random Forest", "• train_model()", "• predict_price()", "• Model auto-load", "• MAE reporting"]),
    (Inches(8.8),  Inches(2.5),  Inches(2.4), Inches(3.5), PURPLE, "DATA LAYER\n(CSV + JSON)",
     ["• prices_3year.csv", "• production.csv", "• users.json", "• crops.json", "• districts.csv"]),
    (Inches(11.6), Inches(2.5),  Inches(1.3), Inches(3.5), RED,    "AUTH\n(JWT)",
     ["• bcrypt hash", "• Token issue", "• Farmer/Consumer", "• Separate roles"]),
]

for (cx, cy, cw, ch, col, title, items) in boxes:
    add_card(slide, cx, cy, cw, ch, col)
    add_rect(slide, cx, cy, cw, Inches(0.08), col)
    # Title
    for j, line in enumerate(title.split("\n")):
        bold = (j == 0)
        fcol = col if j == 0 else TEXT_SUB
        fsz  = Pt(13) if j == 0 else Pt(10)
        add_text(slide, line, cx + Inches(0.1), cy + Inches(0.1 + j * 0.4),
                 cw - Inches(0.15), Inches(0.38),
                 font_size=fsz, bold=bold, color=fcol, align=PP_ALIGN.CENTER)
    for k, item in enumerate(items):
        add_text(slide, item,
                 cx + Inches(0.1), cy + Inches(1.2 + k * 0.44),
                 cw - Inches(0.15), Inches(0.38),
                 font_size=Pt(9.5), color=TEXT_SUB)

# Arrows between boxes (using thin rectangles as connectors)
arrow_y = Inches(4.1)
arrow_positions = [Inches(2.82), Inches(5.62), Inches(8.42), Inches(11.22)]
arrow_colors    = [BLUE, GREEN, AMBER, PURPLE]
for ax, ac in zip(arrow_positions, arrow_colors):
    add_rect(slide, ax, arrow_y, Inches(0.35), Inches(0.08), ac)
    # Arrowhead triangle
    add_rect(slide, ax + Inches(0.28), arrow_y - Inches(0.06), Inches(0.1), Inches(0.2), ac)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — API Endpoints
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "REST API Reference",
             "FastAPI — Auto-documented at http://localhost:8000/docs", GREEN)

endpoints = [
    ("AUTH",      GREEN,  [
        ("POST", "/auth/register",       "Register new farmer or consumer"),
        ("POST", "/auth/login",          "JWT token login"),
        ("GET",  "/auth/profile",        "Get current user profile"),
    ]),
    ("CROP DATA", BLUE,   [
        ("GET",  "/api/vegetables",                              "Search vegetables by name/district"),
        ("GET",  "/api/vegetables-list",                         "Return all known vegetable names"),
        ("GET",  "/api/districts",                               "Return all 13 AP districts with GPS"),
        ("GET",  "/api/predict-price?vegetable=Tomato&district=Guntur&month=6&year=2025",
                                                                 "AI price prediction"),
        ("GET",  "/api/recommendations?district=Guntur",         "Top 8 crop recommendations"),
        ("GET",  "/api/model-info",                              "Model MAE & training stats"),
    ]),
    ("FARMER",    AMBER,  [
        ("POST", "/farmer/list-crop",    "Add new crop listing"),
        ("GET",  "/farmer/my-crops",     "Get logged-in farmer's crops"),
        ("DELETE", "/farmer/crop/{id}", "Remove a listing"),
    ]),
    ("AI / NLP",  PURPLE, [
        ("POST", "/api/voice-command",   "Parse voice text, return intent & speech response"),
        ("POST", "/api/train-model",     "Trigger model re-train (admin)"),
    ]),
]

col_x  = [Inches(0.35), Inches(6.85)]
col_y  = [Inches(1.5),  Inches(1.5)]
group_counter = [0, 0]
current_col   = 0

for group_title, col, eps in endpoints:
    cx = col_x[current_col]
    cy_base = col_y[current_col]
    cw = Inches(6.25)
    ch = Inches(0.48 + len(eps) * 0.72)
    add_card(slide, cx, cy_base, cw, ch, col)
    add_rect(slide, cx, cy_base, cw, Inches(0.44), col)
    add_text(slide, group_title,
             cx + Inches(0.15), cy_base + Inches(0.05),
             cw - Inches(0.2), Inches(0.38),
             font_size=Pt(12), bold=True, color=WHITE)
    for k, (method, path, desc) in enumerate(eps):
        ey = cy_base + Inches(0.52 + k * 0.72)
        m_color = GREEN if method == "GET" else (BLUE if method == "POST" else RED)
        add_rect(slide, cx + Inches(0.12), ey + Inches(0.06),
                 Inches(0.6), Inches(0.32), m_color)
        add_text(slide, method,
                 cx + Inches(0.12), ey + Inches(0.04),
                 Inches(0.6), Inches(0.32),
                 font_size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, path,
                 cx + Inches(0.78), ey,
                 cw - Inches(0.85), Inches(0.32),
                 font_size=Pt(9.5), bold=True, color=col)
        add_text(slide, desc,
                 cx + Inches(0.78), ey + Inches(0.3),
                 cw - Inches(0.85), Inches(0.3),
                 font_size=Pt(9), color=TEXT_SUB)
    col_y[current_col] += ch + Inches(0.18)
    current_col = (current_col + 1) % 2


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Impact & Statistics
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "Impact & Coverage",
             "Real-World Agricultural Transparency for Andhra Pradesh", GREEN)

big_stats = [
    (GREEN,  "13",       "Districts Covered",       "All 13 AP districts\nwith GPS coordinates"),
    (BLUE,   "200+",     "Registered Farmers",       "With crop listings\n& direct contacts"),
    (AMBER,  "3 Years",  "Historical Price Data",    "For accurate ML\ntraining & predictions"),
    (PURPLE, "15+",      "Vegetables Tracked",       "Including Tomato, Onion,\nBrinjal, Chilli & more"),
]

stat_x = Inches(0.5)
for col, number, label, sub in big_stats:
    add_card(slide, stat_x, Inches(1.55), Inches(2.9), Inches(2.6), col)
    add_rect(slide, stat_x, Inches(1.55), Inches(2.9), Inches(0.07), col)
    add_text(slide, number,
             stat_x + Inches(0.1), Inches(1.75),
             Inches(2.7), Inches(0.85),
             font_size=Pt(42), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             stat_x + Inches(0.1), Inches(2.6),
             Inches(2.7), Inches(0.4),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sub,
             stat_x + Inches(0.1), Inches(3.0),
             Inches(2.7), Inches(0.6),
             font_size=Pt(10), color=TEXT_SUB, align=PP_ALIGN.CENTER)
    stat_x += Inches(3.05)

# Benefits table
add_card(slide, Inches(0.35), Inches(4.3), Inches(12.63), Inches(2.9), GREEN)
add_rect(slide, Inches(0.35), Inches(4.3), Inches(12.63), Inches(0.5), GREEN)
add_text(slide, "PLATFORM BENEFITS",
         Inches(0.5), Inches(4.35),
         Inches(5), Inches(0.43),
         font_size=Pt(12.5), bold=True, color=WHITE)

benefits = [
    ("For Farmers 🌾",
     "• Direct access to consumers — eliminate middlemen\n• AI predicted prices help plan harvest timing\n• Voice commands for easy listing (vernacular friendly)"),
    ("For Consumers 🛒",
     "• Find fresh produce near them on the map\n• Price transparency — know market rates before buying\n• Direct farmer contact — call, negotiate, collect"),
    ("For Andhra Pradesh 🏛️",
     "• Digital transparency in agricultural pricing\n• Reduced food waste via efficient supply matching\n• Data-driven policy insights via AI model metrics"),
]

bx = Inches(0.5)
for title, body in benefits:
    add_text(slide, title,
             bx, Inches(4.88),
             Inches(4.1), Inches(0.38),
             font_size=Pt(12), bold=True, color=GREEN_LIGHT)
    add_text(slide, body,
             bx, Inches(5.28),
             Inches(4.1), Inches(1.6),
             font_size=Pt(10), color=TEXT_SUB)
    bx += Inches(4.3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Getting Started / Deployment
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide_header(slide, "Getting Started",
             "Local Setup — 2 Commands. Run Both Services & Open Browser.", BLUE)

steps_deploy = [
    (GREEN,  "STEP 1", "Run Backend",
     "Double-click or run:\n  start_backend.bat\n\nThis installs Python deps (pip install -r requirements.txt)\nand starts FastAPI on http://localhost:8000\n\n✓ Swagger UI available at: http://localhost:8000/docs"),
    (BLUE,   "STEP 2", "Run Frontend",
     "Double-click or run:\n  start_frontend.bat\n\nThis installs Node.js deps (npm install)\nand starts Vite dev server on http://localhost:5173\n\n✓ React app auto-opens in browser"),
    (AMBER,  "STEP 3", "Open Platform",
     "Navigate to:\n  http://localhost:5173\n\nChoose Farmer Portal or Consumer Portal.\nLogin or register a new account.\n\n✓ ML model trains automatically on first launch"),
]

for i, (col, step, title, body) in enumerate(steps_deploy):
    cx = Inches(0.35 + i * 4.35)
    cy = Inches(1.55)
    cw = Inches(4.1)
    ch = Inches(5.65)
    add_card(slide, cx, cy, cw, ch, col)
    add_rect(slide, cx, cy, cw, Inches(0.5), col)
    add_text(slide, step,
             cx + Inches(0.15), cy + Inches(0.07),
             Inches(1.2), Inches(0.38),
             font_size=Pt(11), bold=True, color=WHITE)
    add_text(slide, title,
             cx + Inches(0.15), cy + Inches(0.62),
             cw - Inches(0.25), Inches(0.38),
             font_size=Pt(16), bold=True, color=col)
    add_text(slide, body,
             cx + Inches(0.2), cy + Inches(1.1),
             cw - Inches(0.3), ch - Inches(1.3),
             font_size=Pt(11), color=TEXT_SUB)

# Prerequisites note
add_card(slide, Inches(0.35), Inches(7.1), Inches(12.63), Inches(0.32), PURPLE)
add_text(slide,
         "Prerequisites:  Python 3.8+  ·  Node.js 18+  ·  Git (optional)",
         Inches(0.5), Inches(7.12),
         Inches(12.4), Inches(0.28),
         font_size=Pt(11), color=PURPLE, align=PP_ALIGN.CENTER, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Thank You / Close
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), GREEN)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), GREEN)

add_text(slide, "🌿",
         Inches(5.5), Inches(1.0),
         Inches(2.3), Inches(1.4),
         font_size=Pt(72), align=PP_ALIGN.CENTER)

add_text(slide, "Thank You",
         Inches(2.0), Inches(2.45),
         Inches(9.3), Inches(1.0),
         font_size=Pt(56), bold=True, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(3.5), Inches(3.5), Inches(6.3), GREEN)

add_text(slide, "AgriSmart — Empowering Andhra Pradesh Farmers with AI",
         Inches(2.0), Inches(3.65),
         Inches(9.3), Inches(0.55),
         font_size=Pt(17), color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Digital Transparency Initiative · AI Price Prediction · Direct Farmer-Consumer Marketplace",
         Inches(2.0), Inches(4.25),
         Inches(9.3), Inches(0.45),
         font_size=Pt(12), color=TEXT_SUB, align=PP_ALIGN.CENTER, italic=True)

links = [
    ("Frontend", "http://localhost:5173", BLUE_LIGHT),
    ("API Docs",  "http://localhost:8000/docs", GREEN_LIGHT),
    ("Health",   "http://localhost:8000/health", AMBER),
]
link_x = Inches(2.0)
for lbl, url, col in links:
    add_rect(slide, link_x, Inches(5.1), Inches(3.0), Inches(0.55), CARD_BG,
             line_color=col, line_width=Pt(1))
    add_text(slide, lbl,
             link_x + Inches(0.1), Inches(5.12),
             Inches(2.8), Inches(0.26),
             font_size=Pt(10), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, url,
             link_x + Inches(0.1), Inches(5.37),
             Inches(2.8), Inches(0.25),
             font_size=Pt(9), color=TEXT_SUB, align=PP_ALIGN.CENTER)
    link_x += Inches(3.1)


# ── Save ─────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), "AgriSmart_Presentation.pptx")
prs.save(OUT)
print(f"✅  Saved: {OUT}")
print(f"    Slides: {len(prs.slides)}")
