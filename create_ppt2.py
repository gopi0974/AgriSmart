"""
AgriSmart - Rich Interactive PPT with Images & Effects
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, glob

# ── Paths ─────────────────────────────────────────────────────────────────────
BRAIN = r"C:\Users\GOPI\.gemini\antigravity\brain\f5a1c93d-b553-4ceb-a464-a334d712234f"
def img(name):
    matches = glob.glob(os.path.join(BRAIN, f"{name}_*.png"))
    return matches[0] if matches else None

IMG_HERO     = img("agrismart_hero_farm")
IMG_PROBLEM  = img("problem_farmer_exploitation")
IMG_AI       = img("ai_ml_engine")
IMG_MAP      = img("ap_live_map")
IMG_FARMER   = img("happy_farmer")
IMG_CONSUMER = img("consumer_market")
IMG_TECH     = img("tech_stack_visual")
IMG_VEGGIES  = img("vegetable_flatlay")

# ── Colors ────────────────────────────────────────────────────────────────────
G   = RGBColor(0x16,0xa3,0x4a)
GL  = RGBColor(0x4a,0xde,0x80)
B   = RGBColor(0x25,0x63,0xeb)
BL  = RGBColor(0x93,0xc5,0xfd)
A   = RGBColor(0xf5,0x9e,0x0b)
P   = RGBColor(0xa8,0x55,0xf7)
R   = RGBColor(0xef,0x44,0x44)
BG  = RGBColor(0x0a,0x0f,0x1e)
C1  = RGBColor(0x0f,0x2a,0x1a)
C2  = RGBColor(0x0a,0x1a,0x30)
W   = RGBColor(0xFF,0xFF,0xFF)
TS  = RGBColor(0x94,0xa3,0xb8)
TM  = RGBColor(0xe2,0xe8,0xf0)

SW, SH = Inches(13.33), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, lc=None, lw=Pt(1.5)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc; s.line.width = lw
    else: s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, sz=Pt(13), bold=False, color=TM,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def img_add(slide, path, x, y, w, h):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)

def add_oval(slide, x, y, w, h, fill, lc=None, lw=Pt(1)):
    s = slide.shapes.add_shape(9, x, y, w, h)  # 9 = oval
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc; s.line.width = lw
    else: s.line.fill.background()
    return s

def line(slide, x, y, w, color=G, h=Pt(3)):
    r = rect(slide, x, y, w, h, color); return r

def add_trans(slide, effect="fade"):
    """Add slide transition XML"""
    spTree = slide.shapes._spTree
    mc = slide._element
    # Remove old transition
    for t in mc.findall(qn('p:transition')):
        mc.remove(t)
    if effect == "fade":
        t_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    elif effect == "push":
        t_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>'
    elif effect == "wipe":
        t_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="l"/></p:transition>'
    elif effect == "zoom":
        t_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:zoom/></p:transition>'
    else:
        t_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    mc.append(etree.fromstring(t_xml))

def add_entrance_anim(shape, delay_ms=0):
    """Add fly-in animation to a shape"""
    sp = shape._element
    slide_el = sp.getparent().getparent()
    timing = slide_el.find(qn('p:timing'))
    if timing is None:
        timing_xml = '''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
  <p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">
  <p:childTnLst></p:childTnLst></p:cTn><p:prevCondLst><p:cond evt="onPrevClick" delay="0"><p:tn><p:tgtEl><p:sldTgt/></p:tgtEl></p:tn></p:cond></p:prevCondLst>
  <p:nextCondLst><p:cond evt="onNextClick" delay="0"><p:tn><p:tgtEl><p:sldTgt/></p:tgtEl></p:tn></p:cond></p:nextCondLst>
  </p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>
  <p:bldLst/></p:timing>'''
        slide_el.append(etree.fromstring(timing_xml))
        timing = slide_el.find(qn('p:timing'))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO / COVER  (full bleed image + overlay)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, BG)
add_trans(s, "zoom")

# Full slide hero image (left 55%)
img_add(s, IMG_HERO, 0, 0, Inches(7.5), SH)

# Dark overlay gradient on the image
r1 = rect(s, 0, 0, Inches(7.5), SH, BG)
r1.fill.solid(); r1.fill.fore_color.rgb = BG
r1._element.spPr.solidFill.srgbClr.set('val', '0a0f1e')
# Just use a semi-transparent look via shape transparency trick
# Actually let's do a diagonal overlay shape - use a darker green approach
rect(s, 0, 0, Inches(4.0), SH, RGBColor(0x0a,0x0f,0x1e))

# Right panel - dark content area
rect(s, Inches(6.5), 0, Inches(6.83), SH, RGBColor(0x05,0x0a,0x18))

# Accent vertical bar
rect(s, Inches(6.3), 0, Inches(0.12), SH, G)

# Keep image visible on right portion - overlap with image
img_add(s, IMG_VEGGIES, Inches(7.5), Inches(2.5), Inches(5.83), Inches(5.0))
rect(s, Inches(7.5), Inches(2.5), Inches(5.83), Inches(5.0), RGBColor(0x05,0x0a,0x18))

# Glow circle decoration
add_oval(s, Inches(0.5), Inches(1.0), Inches(3.0), Inches(3.0), RGBColor(0x0d,0x2a,0x18))

# Logo / badge
rect(s, Inches(6.6), Inches(0.35), Inches(0.5), Inches(0.5), G)
add_oval(s, Inches(6.58), Inches(0.33), Inches(0.54), Inches(0.54), G)
txt(s, "🌿", Inches(6.6), Inches(0.3), Inches(0.55), Inches(0.6), sz=Pt(22), align=PP_ALIGN.CENTER)

txt(s, "AgriSmart", Inches(6.65), Inches(1.0), Inches(6.5), Inches(1.4),
    sz=Pt(58), bold=True, color=GL, align=PP_ALIGN.LEFT)
line(s, Inches(6.65), Inches(2.4), Inches(5.5), G, Pt(4))

txt(s, "AI-Powered Vegetable Marketplace", Inches(6.65), Inches(2.55), Inches(6.5), Inches(0.6),
    sz=Pt(18), bold=True, color=W)
txt(s, "Connecting Farmers & Consumers\nacross Andhra Pradesh with\nReal-time AI Price Predictions",
    Inches(6.65), Inches(3.2), Inches(6.2), Inches(1.0), sz=Pt(13), color=TS)

# Stat pills (horizontal)
stats = [("13\nDistricts", G), ("200+\nFarmers", B), ("3yr\nData", A), ("AI\nPrices", P)]
px = Inches(6.65)
for val, c in stats:
    add_oval(s, px, Inches(4.4), Inches(1.4), Inches(1.1), C2, lc=c, lw=Pt(2))
    txt(s, val, px, Inches(4.45), Inches(1.4), Inches(1.0), sz=Pt(11), bold=True, color=c, align=PP_ALIGN.CENTER)
    px += Inches(1.5)

txt(s, "Digital Transparency Initiative · Andhra Pradesh",
    Inches(6.65), Inches(6.9), Inches(6.4), Inches(0.4), sz=Pt(9), color=TS, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, RGBColor(0x10,0x05,0x05))
add_trans(s, "wipe")

# Full-width mood image (right half)
img_add(s, IMG_PROBLEM, Inches(5.5), 0, Inches(7.83), SH)
rect(s, Inches(5.5), 0, Inches(7.83), SH, RGBColor(0x15,0x05,0x05))  # tint overlay

# Dark left column
rect(s, 0, 0, Inches(6.0), SH, RGBColor(0x0d,0x03,0x03))
rect(s, Inches(5.9), 0, Inches(0.14), SH, R)

# Header
txt(s, "⚠ THE PROBLEM", Inches(0.4), Inches(0.25), Inches(5.5), Inches(0.45),
    sz=Pt(11), bold=True, color=R)
txt(s, "Why Farmers\nGet Exploited", Inches(0.4), Inches(0.68), Inches(5.4), Inches(1.0),
    sz=Pt(30), bold=True, color=W)
line(s, Inches(0.4), Inches(1.68), Inches(4.8), R, Pt(3))

problems = [
    ("💸 Price Opacity",      "Farmers get only 20-30% of final retail price"),
    ("🔗 Middlemen Chain",    "3–5 middlemen between farm & consumer"),
    ("📍 No Direct Channel",  "No platform to connect farmer ↔ consumer"),
    ("📉 Price Volatility",   "Seasonal swings with no advance warning"),
    ("🗺 Info Gap",           "Consumers can't locate fresh produce nearby"),
]
for i, (title, desc) in enumerate(problems):
    cy = Inches(1.9 + i * 1.0)
    rect(s, Inches(0.35), cy, Inches(5.4), Inches(0.88), RGBColor(0x20,0x08,0x08), lc=R, lw=Pt(1))
    rect(s, Inches(0.35), cy, Inches(0.09), Inches(0.88), R)
    txt(s, title, Inches(0.55), cy+Inches(0.06), Inches(5.0), Inches(0.33), sz=Pt(12), bold=True, color=RGBColor(0xff,0x88,0x88))
    txt(s, desc,  Inches(0.55), cy+Inches(0.39), Inches(5.0), Inches(0.35), sz=Pt(10.5), color=TS)

# Right - stat callout
rect(s, Inches(7.2), Inches(1.5), Inches(3.8), Inches(1.3), RGBColor(0x2a,0x08,0x08), lc=R, lw=Pt(2))
txt(s, "Only 20-30%", Inches(7.3), Inches(1.6), Inches(3.6), Inches(0.6), sz=Pt(28), bold=True, color=R, align=PP_ALIGN.CENTER)
txt(s, "of retail price reaches the farmer", Inches(7.3), Inches(2.1), Inches(3.6), Inches(0.4), sz=Pt(11), color=TS, align=PP_ALIGN.CENTER)

rect(s, Inches(7.2), Inches(3.1), Inches(3.8), Inches(1.3), RGBColor(0x1a,0x08,0x08), lc=A, lw=Pt(2))
txt(s, "3–5 Middlemen", Inches(7.3), Inches(3.2), Inches(3.6), Inches(0.6), sz=Pt(26), bold=True, color=A, align=PP_ALIGN.CENTER)
txt(s, "between farm gate and consumer", Inches(7.3), Inches(3.7), Inches(3.6), Inches(0.4), sz=Pt(11), color=TS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION (split layout with vegetable image)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, BG)
add_trans(s, "push")

# Left image panel
img_add(s, IMG_VEGGIES, 0, 0, Inches(5.5), SH)
rect(s, Inches(4.2), 0, Inches(1.3), SH, BG)  # fade edge

# Right content
rect(s, Inches(5.0), 0, Inches(8.33), SH, RGBColor(0x07,0x10,0x0d))
rect(s, Inches(5.0), 0, Inches(0.12), SH, G)

txt(s, "✅ OUR SOLUTION", Inches(5.25), Inches(0.28), Inches(7.8), Inches(0.4),
    sz=Pt(11), bold=True, color=GL)
txt(s, "AgriSmart Platform", Inches(5.25), Inches(0.68), Inches(7.8), Inches(0.7),
    sz=Pt(32), bold=True, color=W)
line(s, Inches(5.25), Inches(1.4), Inches(7.6), G, Pt(3))

solutions = [
    (G,  "🌾", "Farmer Portal",      "Register, list crops & quantities, get AI price guidance, reach thousands of consumers directly."),
    (B,  "🛒", "Consumer Portal",    "Search nearby farmers on live map, compare AI-predicted prices, contact farmers directly."),
    (A,  "🤖", "AI Price Engine",    "Random Forest ML model on 3-year AP market data. Real-time ₹/kg predictions per district."),
    (P,  "🗺️", "Live Location Map",  "Leaflet.js map showing every farmer pin with contact info, vegetable type & quantity."),
]
for i, (c, icon, title, body) in enumerate(solutions):
    cx = Inches(5.3 if i%2==0 else 9.3)
    cy = Inches(1.6 + (i//2)*2.55)
    rect(s, cx, cy, Inches(3.75), Inches(2.3), RGBColor(0x0c,0x1e,0x18), lc=c, lw=Pt(1.5))
    rect(s, cx, cy, Inches(3.75), Inches(0.07), c)
    txt(s, f"{icon} {title}", cx+Inches(0.18), cy+Inches(0.14), Inches(3.4), Inches(0.4),
        sz=Pt(13), bold=True, color=c)
    txt(s, body, cx+Inches(0.18), cy+Inches(0.6), Inches(3.4), Inches(1.5),
        sz=Pt(10.5), color=TS)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FARMER PORTAL (happy farmer image)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, RGBColor(0x04,0x12,0x08))
add_trans(s, "fade")

# Top green bar
rect(s, 0, 0, SW, Inches(0.08), G)
# Bottom green bar
rect(s, 0, SH-Inches(0.08), SW, Inches(0.08), G)

# Left content area
txt(s, "🌾 FARMER PORTAL", Inches(0.5), Inches(0.2), Inches(6.5), Inches(0.45),
    sz=Pt(11), bold=True, color=GL)
txt(s, "Grow. List.\nEarn More.", Inches(0.5), Inches(0.65), Inches(6.5), Inches(1.4),
    sz=Pt(38), bold=True, color=W)
line(s, Inches(0.5), Inches(2.05), Inches(5.8), G, Pt(4))
txt(s, "Farmers register on AgriSmart, list their harvest\nwith quantity & district, and get AI price insights —\nall without any middleman.",
    Inches(0.5), Inches(2.2), Inches(6.2), Inches(1.0), sz=Pt(13), color=TS)

# Feature list
features = [
    ("📋 Simple Registration", "Name, mobile, district — get started in 60s"),
    ("🥦 Crop Listing",        "Add vegetable, kg quantity, GPS location"),
    ("💰 AI Price Guide",      "See predicted market price before listing"),
    ("🎤 Voice Commands",      "Say 'Add 100kg Tomato in Guntur' to auto-fill"),
    ("📞 Direct Contact",      "Consumers see your number & call directly"),
]
for i, (t, d) in enumerate(features):
    cy = Inches(3.35 + i * 0.74)
    rect(s, Inches(0.45), cy, Inches(6.3), Inches(0.65), RGBColor(0x0a,0x22,0x10), lc=G, lw=Pt(1))
    add_oval(s, Inches(0.5), cy+Inches(0.15), Inches(0.35), Inches(0.35), G)
    txt(s, "✓", Inches(0.5), cy+Inches(0.13), Inches(0.37), Inches(0.35), sz=Pt(10), bold=True, color=W, align=PP_ALIGN.CENTER)
    txt(s, t, Inches(1.0), cy+Inches(0.04), Inches(3.5), Inches(0.28), sz=Pt(11), bold=True, color=GL)
    txt(s, d, Inches(1.0), cy+Inches(0.32), Inches(5.5), Inches(0.28), sz=Pt(10), color=TS)

# Right - farmer image in circle frame
img_add(s, IMG_FARMER, Inches(7.0), Inches(0.5), Inches(5.8), Inches(6.7))
# Circular frame overlay (dark corners via ovals)
rect(s, Inches(7.0), Inches(0.5), Inches(5.8), Inches(6.7), RGBColor(0x04,0x12,0x08))
img_add(s, IMG_FARMER, Inches(7.2), Inches(0.6), Inches(5.5), Inches(6.5))
add_oval(s, Inches(7.1), Inches(0.45), Inches(5.8), Inches(6.7), RGBColor(0x0a,0x22,0x12), lc=G, lw=Pt(3))
img_add(s, IMG_FARMER, Inches(7.25), Inches(0.65), Inches(5.5), Inches(6.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONSUMER PORTAL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, RGBColor(0x04,0x09,0x18))
add_trans(s, "push")

rect(s, 0, 0, SW, Inches(0.08), B)
rect(s, 0, SH-Inches(0.08), SW, Inches(0.08), B)

# Left - consumer image
img_add(s, IMG_CONSUMER, Inches(0.3), Inches(0.5), Inches(5.8), Inches(6.5))
add_oval(s, Inches(0.2), Inches(0.45), Inches(5.9), Inches(6.6), RGBColor(0x06,0x12,0x28), lc=B, lw=Pt(3))
img_add(s, IMG_CONSUMER, Inches(0.35), Inches(0.6), Inches(5.7), Inches(6.3))

# Right content
txt(s, "🛒 CONSUMER PORTAL", Inches(6.5), Inches(0.2), Inches(6.5), Inches(0.45),
    sz=Pt(11), bold=True, color=BL)
txt(s, "Find. Compare.\nBuy Fresh.", Inches(6.5), Inches(0.65), Inches(6.5), Inches(1.4),
    sz=Pt(38), bold=True, color=W)
line(s, Inches(6.5), Inches(2.05), Inches(6.4), B, Pt(4))
txt(s, "Consumers find fresh vegetables from nearby\nfarmers, check AI price predictions, and\ncontact farmers directly — no middlemen.",
    Inches(6.5), Inches(2.2), Inches(6.5), Inches(1.0), sz=Pt(13), color=TS)

features = [
    ("🔍 Smart Search",        "Search by vegetable name or district"),
    ("🗺️ Map View",            "Pin-drop map of all farmers near you"),
    ("📊 Price Predictions",   "AI-predicted ₹/kg before you negotiate"),
    ("📞 Direct Contact",      "Farmer mobile number shown on card"),
    ("⭐ Recommendations",     "Top 8 best-value crops in your district"),
]
for i, (t, d) in enumerate(features):
    cy = Inches(3.35 + i * 0.74)
    rect(s, Inches(6.5), cy, Inches(6.5), Inches(0.65), RGBColor(0x08,0x14,0x28), lc=B, lw=Pt(1))
    add_oval(s, Inches(6.55), cy+Inches(0.15), Inches(0.35), Inches(0.35), B)
    txt(s, "✓", Inches(6.55), cy+Inches(0.13), Inches(0.37), Inches(0.35), sz=Pt(10), bold=True, color=W, align=PP_ALIGN.CENTER)
    txt(s, t, Inches(7.05), cy+Inches(0.04), Inches(3.5), Inches(0.28), sz=Pt(11), bold=True, color=BL)
    txt(s, d, Inches(7.05), cy+Inches(0.32), Inches(5.8), Inches(0.28), sz=Pt(10), color=TS)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE MAP
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, RGBColor(0x04,0x09,0x18))
add_trans(s, "wipe")

# Full bleed map image
img_add(s, IMG_MAP, 0, 0, SW, SH)
# Dark overlay
rect(s, 0, 0, SW, SH, RGBColor(0x04,0x09,0x18))
img_add(s, IMG_MAP, 0, 0, SW, SH)

# Top overlay panel
rect(s, 0, 0, SW, Inches(1.5), RGBColor(0x04,0x09,0x18))
rect(s, 0, 0, SW, Inches(0.07), B)

txt(s, "🗺️ LIVE INTERACTIVE MAP", Inches(0.5), Inches(0.18), Inches(8.0), Inches(0.45),
    sz=Pt(14), bold=True, color=BL)
txt(s, "Real-time farmer location mapping across all 13 Andhra Pradesh districts",
    Inches(0.5), Inches(0.62), Inches(10.5), Inches(0.45), sz=Pt(13), color=TS)
line(s, Inches(0.5), Inches(1.05), Inches(12.3), B, Pt(2))

# Bottom legend panel
rect(s, 0, Inches(6.4), SW, Inches(1.1), RGBColor(0x04,0x09,0x18))
line(s, Inches(0.5), Inches(6.42), Inches(12.3), B, Pt(2))

legend = [
    (G, "🟢 Live Farmer Pins",   "GPS-accurate drop pins for each registered farmer"),
    (B, "🔵 Market Points",      "Historical market average price per district"),
    (A, "🟡 Your Location",      "Consumer's current GPS position"),
    (P, "🟣 Recommendations",    "Top 8 AI-recommended crops in your district"),
]
lx = Inches(0.5)
for c, icon, desc in legend:
    rect(s, lx, Inches(6.5), Inches(3.1), Inches(0.82), RGBColor(0x07,0x12,0x28), lc=c, lw=Pt(1.5))
    txt(s, icon, lx+Inches(0.1), Inches(6.53), Inches(2.9), Inches(0.3), sz=Pt(10), bold=True, color=c)
    txt(s, desc, lx+Inches(0.1), Inches(6.8), Inches(2.9), Inches(0.28), sz=Pt(9), color=TS)
    lx += Inches(3.25)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI / ML ENGINE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, RGBColor(0x06,0x0e,0x05))
add_trans(s, "zoom")

# AI image top-right corner background
img_add(s, IMG_AI, Inches(5.5), 0, Inches(7.83), SH)
rect(s, Inches(5.5), 0, Inches(3.0), SH, RGBColor(0x06,0x0e,0x05))  # blend edge

rect(s, 0, 0, SW, Inches(0.07), A)
rect(s, 0, SH-Inches(0.08), SW, Inches(0.08), A)

# Left content
txt(s, "🤖 AI PRICE ENGINE", Inches(0.5), Inches(0.2), Inches(6.5), Inches(0.4),
    sz=Pt(11), bold=True, color=A)
txt(s, "Random Forest\nML Model", Inches(0.5), Inches(0.62), Inches(6.5), Inches(1.1),
    sz=Pt(34), bold=True, color=W)
line(s, Inches(0.5), Inches(1.72), Inches(5.5), A, Pt(4))

# Pipeline steps
steps = [
    ("1", A,  "Data Ingestion",     "3-year AP price history (1.8 MB CSV) + yield, soil, weather data"),
    ("2", GL, "Feature Engineering","Variables: vegetable, district, month, year, avg_yield, avg_area"),
    ("3", B,  "Model Training",     "RandomForestRegressor → serialized with joblib for instant reload"),
    ("4", P,  "Auto-load Startup",  "FastAPI checks for saved model; trains automatically if missing"),
    ("5", A,  "Real-time Predict",  "GET /api/predict-price returns ₹/kg within milliseconds"),
]
for i, (num, c, title, body) in enumerate(steps):
    cy = Inches(1.9 + i * 1.05)
    add_oval(s, Inches(0.45), cy+Inches(0.08), Inches(0.52), Inches(0.52), c)
    txt(s, num, Inches(0.45), cy+Inches(0.07), Inches(0.52), Inches(0.48),
        sz=Pt(14), bold=True, color=W, align=PP_ALIGN.CENTER)
    rect(s, Inches(1.1), cy+Inches(0.02), Inches(5.1), Inches(0.92), RGBColor(0x0c,0x1e,0x0a), lc=c, lw=Pt(1))
    txt(s, title, Inches(1.22), cy+Inches(0.06), Inches(4.8), Inches(0.32), sz=Pt(12), bold=True, color=c)
    txt(s, body, Inches(1.22), cy+Inches(0.42), Inches(4.8), Inches(0.42), sz=Pt(10.5), color=TS)
    if i < 4:
        rect(s, Inches(0.68), cy+Inches(0.6), Inches(0.05), Inches(0.5), c)

# Right panel - metrics over image
rect(s, Inches(8.5), Inches(1.2), Inches(4.5), Inches(5.8), RGBColor(0x0a,0x18,0x07))
rect(s, Inches(8.5), Inches(1.2), Inches(4.5), Inches(0.07), A)
txt(s, "MODEL METRICS", Inches(8.65), Inches(1.3), Inches(4.2), Inches(0.4),
    sz=Pt(12), bold=True, color=A)

metrics = [
    ("Algorithm",  "Random Forest Regressor"),
    ("Target",     "avg_price (₹ per kg)"),
    ("Metric",     "Mean Absolute Error (MAE)"),
    ("Training",   "Auto on first launch"),
    ("Inference",  "< 50ms per prediction"),
    ("Coverage",   "15+ vegetables × 13 districts"),
    ("API",        "GET /api/predict-price"),
    ("Retrain",    "POST /api/train-model"),
]
for i, (k, v) in enumerate(metrics):
    cy = Inches(1.8 + i * 0.6)
    rect(s, Inches(8.6), cy, Inches(4.2), Inches(0.52), RGBColor(0x10,0x20,0x0d), lc=GL, lw=Pt(0.5))
    txt(s, k+":", Inches(8.72), cy+Inches(0.08), Inches(1.5), Inches(0.3), sz=Pt(10), bold=True, color=GL)
    txt(s, v, Inches(10.1), cy+Inches(0.08), Inches(2.6), Inches(0.3), sz=Pt(10), color=TS)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, BG)
add_trans(s, "push")

# Background tech image (faint)
img_add(s, IMG_TECH, 0, 0, SW, SH)
rect(s, 0, 0, SW, SH, RGBColor(0x06,0x0d,0x1a))  # heavy overlay

rect(s, 0, 0, SW, Inches(0.07), B)
txt(s, "⚙ TECH STACK", Inches(0.5), Inches(0.18), Inches(5.0), Inches(0.4),
    sz=Pt(11), bold=True, color=BL)
txt(s, "Modern Open-Source Architecture",
    Inches(0.5), Inches(0.6), Inches(9.0), Inches(0.55), sz=Pt(24), bold=True, color=W)
line(s, Inches(0.5), Inches(1.17), Inches(12.3), B, Pt(2))

cols = [
    ("FRONTEND", B, [
        ("React 18 + TypeScript", "Type-safe component UI"),
        ("Vite", "Ultra-fast HMR dev server"),
        ("Tailwind CSS", "Utility-first responsive styles"),
        ("Leaflet.js", "Interactive map with geo-pins"),
        ("Lucide Icons", "Crisp SVG icon library"),
    ]),
    ("BACKEND", G, [
        ("FastAPI", "Async REST API, Swagger docs"),
        ("Python 3.11", "Core application runtime"),
        ("Pandas", "Data wrangling & CSV pipelines"),
        ("Scikit-learn", "Random Forest ML model"),
        ("Joblib", "Model serialization & load"),
    ]),
    ("AUTH & DATA", A, [
        ("JWT + bcrypt", "Secure token authentication"),
        ("users.json", "Lightweight user persistence"),
        ("crops.json", "Live farmer crop listings"),
        ("CSV Datasets", "3yr price, weather, soil, yield"),
        ("districts.csv", "13 AP districts GPS coordinates"),
    ]),
]
for i, (title, c, items) in enumerate(cols):
    cx = Inches(0.4 + i * 4.3)
    cy = Inches(1.35)
    cw, ch = Inches(4.0), Inches(5.9)
    rect(s, cx, cy, cw, ch, RGBColor(0x08,0x12,0x22), lc=c, lw=Pt(2))
    rect(s, cx, cy, cw, Inches(0.52), c)
    txt(s, title, cx+Inches(0.15), cy+Inches(0.08), cw-Inches(0.2), Inches(0.38),
        sz=Pt(13), bold=True, color=W, align=PP_ALIGN.CENTER)
    for j, (tech, desc) in enumerate(items):
        iy = cy + Inches(0.65 + j * 1.02)
        rect(s, cx+Inches(0.15), iy, cw-Inches(0.28), Inches(0.9), RGBColor(0x0c,0x18,0x2c), lc=c, lw=Pt(0.75))
        rect(s, cx+Inches(0.15), iy, Inches(0.06), Inches(0.9), c)
        txt(s, tech, cx+Inches(0.3), iy+Inches(0.06), cw-Inches(0.4), Inches(0.32), sz=Pt(11.5), bold=True, color=c)
        txt(s, desc, cx+Inches(0.3), iy+Inches(0.42), cw-Inches(0.4), Inches(0.36), sz=Pt(9.5), color=TS)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, RGBColor(0x04,0x10,0x08))
add_trans(s, "fade")

# Background farm image (faint top)
img_add(s, IMG_HERO, 0, 0, SW, Inches(3.8))
rect(s, 0, 0, SW, Inches(3.8), RGBColor(0x04,0x10,0x08))
img_add(s, IMG_HERO, 0, 0, SW, Inches(3.5))

rect(s, 0, 0, SW, Inches(0.07), G)
txt(s, "📈 IMPACT & COVERAGE", Inches(0.5), Inches(0.18), Inches(8.0), Inches(0.4),
    sz=Pt(11), bold=True, color=GL)
txt(s, "Andhra Pradesh Agricultural Transformation",
    Inches(0.5), Inches(0.62), Inches(12.0), Inches(0.55), sz=Pt(24), bold=True, color=W)
line(s, Inches(0.5), Inches(1.2), Inches(12.3), G, Pt(2))

# Big stat circles
bigs = [
    (G,  "13",      "Districts"),
    (B,  "200+",    "Farmers"),
    (A,  "3 Years", "Market Data"),
    (P,  "15+",     "Vegetables"),
    (GL, "₹/kg",    "AI Predictions"),
]
sx = Inches(0.5)
for c, val, label in bigs:
    add_oval(s, sx, Inches(1.45), Inches(2.3), Inches(2.0), RGBColor(0x08,0x1e,0x10), lc=c, lw=Pt(2.5))
    txt(s, val,   sx, Inches(1.65), Inches(2.3), Inches(0.85), sz=Pt(28), bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(s, label, sx, Inches(2.5),  Inches(2.3), Inches(0.38), sz=Pt(11), bold=True, color=W,  align=PP_ALIGN.CENTER)
    sx += Inches(2.52)

# Bottom benefit panels
rect(s, 0, Inches(3.65), SW, Inches(3.75), RGBColor(0x04,0x10,0x08))
line(s, Inches(0.5), Inches(3.72), Inches(12.3), G, Pt(1.5))

benefits = [
    (G,  "For Farmers 🌾",      ["Eliminate middlemen — sell direct", "AI price guide = better harvest planning", "Voice listing in Telugu/English", "GPS pin shows buyers your exact location"]),
    (B,  "For Consumers 🛒",    ["Find fresh produce in your district now", "Compare AI prices before negotiating", "Direct farmer mobile number visible", "Map view — see who's closest to you"]),
    (P,  "For AP State 🏛️",     ["Price transparency in agriculture", "Reduced food waste via supply matching", "Data for government policy planning", "Digital empowerment of rural farmers"]),
]
bx = Inches(0.4)
for c, title, pts in benefits:
    rect(s, bx, Inches(3.85), Inches(4.1), Inches(3.45), RGBColor(0x08,0x18,0x10), lc=c, lw=Pt(1.5))
    rect(s, bx, Inches(3.85), Inches(4.1), Inches(0.07), c)
    txt(s, title, bx+Inches(0.18), Inches(3.95), Inches(3.75), Inches(0.38), sz=Pt(13), bold=True, color=c)
    for k, pt in enumerate(pts):
        txt(s, f"  ▸ {pt}", bx+Inches(0.18), Inches(4.42)+Inches(k*0.7), Inches(3.75), Inches(0.62), sz=Pt(11), color=TS)
    bx += Inches(4.3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GETTING STARTED
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, BG)
add_trans(s, "push")

img_add(s, IMG_TECH, 0, 0, SW, SH)
rect(s, 0, 0, SW, SH, RGBColor(0x06,0x0d,0x1a))

rect(s, 0, 0, SW, Inches(0.07), B)
txt(s, "🚀 GETTING STARTED", Inches(0.5), Inches(0.18), Inches(8.0), Inches(0.4),
    sz=Pt(11), bold=True, color=BL)
txt(s, "2 Commands. Running in Minutes.",
    Inches(0.5), Inches(0.62), Inches(10.0), Inches(0.55), sz=Pt(26), bold=True, color=W)
line(s, Inches(0.5), Inches(1.22), Inches(12.3), B, Pt(2))

steps3 = [
    (G,  "STEP 1", "▶ start_backend.bat",
     "Installs Python deps\nStarts FastAPI on :8000\n\nSwagger: localhost:8000/docs\nML model auto-trains if missing\n\n✓ pip install -r requirements.txt\n✓ uvicorn app:app --port 8000"),
    (B,  "STEP 2", "▶ start_frontend.bat",
     "Installs Node.js deps\nStarts Vite on :5173\n\nHot-reload dev server\nAuto opens in browser\n\n✓ npm install\n✓ npm run dev"),
    (A,  "STEP 3", "🌐 Open Platform",
     "Navigate to:\nlocalhost:5173\n\nChoose Farmer or Consumer\nRegister / Login\n\n✓ Farmer: list your crops\n✓ Consumer: find fresh produce"),
]
for i, (c, step, cmd, body) in enumerate(steps3):
    cx = Inches(0.4 + i * 4.3)
    rect(s, cx, Inches(1.45), Inches(4.1), Inches(5.85), RGBColor(0x07,0x11,0x22), lc=c, lw=Pt(2))
    # Step badge
    add_oval(s, cx+Inches(0.18), Inches(1.55), Inches(0.65), Inches(0.65), c)
    txt(s, str(i+1), cx+Inches(0.18), Inches(1.55), Inches(0.67), Inches(0.62),
        sz=Pt(18), bold=True, color=W, align=PP_ALIGN.CENTER)
    txt(s, step, cx+Inches(1.0), Inches(1.6), Inches(3.0), Inches(0.4), sz=Pt(10), bold=True, color=TS)
    txt(s, cmd, cx+Inches(0.18), Inches(2.3), Inches(3.75), Inches(0.45),
        sz=Pt(14), bold=True, color=c)
    line(s, cx+Inches(0.18), Inches(2.78), Inches(3.75), c, Pt(2))

    # Code-like body in monospace style box
    rect(s, cx+Inches(0.18), Inches(2.9), Inches(3.75), Inches(4.1), RGBColor(0x03,0x07,0x15), lc=c, lw=Pt(0.75))
    txt(s, body, cx+Inches(0.3), Inches(3.0), Inches(3.55), Inches(3.9), sz=Pt(11), color=TS)

# Prereqs bar
rect(s, Inches(0.35), Inches(7.15), Inches(12.65), Inches(0.28), RGBColor(0x08,0x12,0x28), lc=B, lw=Pt(1))
txt(s, "Prerequisites:  Python 3.8+  ·  Node.js 18+  ·  Windows / Mac / Linux",
    Inches(0.5), Inches(7.17), Inches(12.0), Inches(0.25), sz=Pt(10), color=BL, align=PP_ALIGN.CENTER, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLK)
bg(s, BG)
add_trans(s, "zoom")

# Full hero image background
img_add(s, IMG_HERO, 0, 0, SW, SH)
rect(s, 0, 0, SW, SH, RGBColor(0x04,0x0a,0x10))
img_add(s, IMG_HERO, 0, 0, SW, SH)
rect(s, 0, 0, SW, SH, RGBColor(0x06,0x0e,0x06))

# Large center content
rect(s, Inches(2.0), Inches(1.0), Inches(9.33), Inches(5.5), RGBColor(0x05,0x12,0x08))
rect(s, Inches(2.0), Inches(1.0), Inches(9.33), Inches(0.09), G)
rect(s, Inches(2.0), Inches(6.41), Inches(9.33), Inches(0.09), G)
rect(s, Inches(2.0), Inches(1.0), Inches(0.07), Inches(5.5), G)
rect(s, Inches(11.26), Inches(1.0), Inches(0.07), Inches(5.5), G)

txt(s, "🌿", Inches(5.5), Inches(1.3), Inches(2.3), Inches(1.1), sz=Pt(52), align=PP_ALIGN.CENTER)
txt(s, "Thank You!", Inches(2.0), Inches(2.4), Inches(9.33), Inches(0.9),
    sz=Pt(50), bold=True, color=GL, align=PP_ALIGN.CENTER)
line(s, Inches(3.5), Inches(3.35), Inches(6.3), G, Pt(4))
txt(s, "AgriSmart — Empowering Andhra Pradesh Farmers with AI",
    Inches(2.0), Inches(3.5), Inches(9.33), Inches(0.55),
    sz=Pt(16), color=W, align=PP_ALIGN.CENTER)
txt(s, "Digital Transparency Initiative · AI Price Prediction · Direct Farmer–Consumer Marketplace",
    Inches(2.0), Inches(4.1), Inches(9.33), Inches(0.4),
    sz=Pt(11), color=TS, align=PP_ALIGN.CENTER, italic=True)

# URL bubbles
links = [("🌐 Frontend", "http://localhost:5173", B),
         ("📖 API Docs", "http://localhost:8000/docs", G),
         ("✅ Health",   "http://localhost:8000/health", A)]
lx = Inches(2.5)
for lbl, url, c in links:
    rect(s, lx, Inches(4.7), Inches(2.6), Inches(0.85), RGBColor(0x08,0x15,0x10), lc=c, lw=Pt(1.5))
    txt(s, lbl, lx+Inches(0.1), Inches(4.75), Inches(2.4), Inches(0.3), sz=Pt(11), bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(s, url,  lx+Inches(0.1), Inches(5.04), Inches(2.4), Inches(0.28), sz=Pt(8.5), color=TS, align=PP_ALIGN.CENTER)
    lx += Inches(2.8)

txt(s, "Created for the Digital Transparency Initiative (DTI)",
    Inches(2.0), Inches(6.1), Inches(9.33), Inches(0.3),
    sz=Pt(9), color=TS, align=PP_ALIGN.CENTER, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AgriSmart_v2.pptx")
prs.save(OUT)
print(f"\n✅ Saved: {OUT}")
print(f"   Slides : {len(prs.slides)}")
print(f"   Images : {sum(1 for f in [IMG_HERO,IMG_PROBLEM,IMG_AI,IMG_MAP,IMG_FARMER,IMG_CONSUMER,IMG_TECH,IMG_VEGGIES] if f)}/8")
